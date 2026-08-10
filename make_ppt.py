import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# PPTX 생성
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 다크모드 컬러 팔레트
BG_DARK = RGBColor(15, 23, 42)
CARD_BG = RGBColor(30, 41, 59)
TEXT_MAIN = RGBColor(241, 245, 249)
TEXT_MUTED = RGBColor(148, 163, 184)
PRIMARY = RGBColor(59, 130, 246)
ACCENT = RGBColor(16, 185, 129)
CARD_BORDER = RGBColor(51, 65, 85)

blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK

def add_header(slide, title_text, category_text="E-COMMERCE DASHBOARD PROJECT"):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.733), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    p0.font.name = "Malgun Gothic"
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p1.font.name = "Malgun Gothic"
    p1.space_before = Pt(3)

def add_card(slide, left, top, width, height, title, items_list):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.font.name = "Malgun Gothic"
        first = False
    else:
        first = True
        
    for item in items_list:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MAIN
        p.font.name = "Malgun Gothic"
        p.space_before = Pt(6)

# Slide 1: 표제
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1)
tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "E-Commerce 주문 데이터 분석 대시보드"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = TEXT_MAIN
p.font.name = "Malgun Gothic"
p2 = tf.add_paragraph()
p2.text = "Streamlit & Plotly 기반 실시간 분석 및 프로젝트 종합 회고 보고서"
p2.font.size = Pt(17)
p2.font.color.rgb = PRIMARY
p2.font.name = "Malgun Gothic"
p2.space_before = Pt(12)

# Slide 2: 프로젝트 개요
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2)
add_header(s2, "1. 프로젝트 개요", "PROJECT OVERVIEW")
add_card(s2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.5), "🎯 프로젝트 목적", [
    "• E-Commerce 쇼핑몰의 주문, 고객, 상품 데이터의 다이내믹 분석",
    "• 사용자 선택 조건(사이드바 필터)에 따른 실시간 KPI/차트 연동",
    "• 데이터 부재(Empty State) 시에도 에러 없이 동작하는 웹 애플리케이션 구축"
])
add_card(s2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5), "🛠️ 개발 환경", [
    "• 언어: Python 3.10+ / 프레임워크: Streamlit",
    "• 데이터 처리: Pandas, NumPy / 시각화: Plotly Express",
    "• 버전 관리: Git, GitHub"
])

# Slide 3: 팀원과 역할
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3)
add_header(s3, "2. 팀원과 역할", "TEAM & ROLES")
add_card(s3, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.3), "👤 팀원 A (팀장)\n[UI/UX 개발]", ["• Streamlit 대시보드 화면 구조 설계", "• 사이드바 동적 필터 연동", "• 핵심 KPI 8종 레이아웃 배치"])
add_card(s3, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.3), "👤 팀원 B\n[데이터 & 시각화]", ["• Pandas 기반 데이터 전처리 및 집계", "• Plotly Express 시각화 차트 2종 구현", "• 차트별 자동 결과 해석 패널 작성"])
add_card(s3, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.3), "👤 팀원 C\n[검증 & 문서화]", ["• AI 코드 기능 단위 검증서 작성", "• Empty State 예외 조건 로직 추가", "• 프로젝트 회고 및 최종 보고서 작성"])

# Slide 4: 핵심 KPI
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4)
add_header(s4, "3. 구현 결과: 데이터셋 & 핵심 KPI", "DASHBOARD IMPLEMENTATION - METRICS")
add_card(s4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "📌 핵심 성과 지표 (8대 KPI)", [
    "• 전체 고객 수: 1,200 명 | 전체 상품 수: 300 개 | 전체 주문 수: 6,000 건",
    "• 총 주문 수량: 23,596 개 | 총 주문 금액: ₩1,837,774,300 원",
    "• 평균 주문 금액: ₩306,296 원 | 배송완료 주문 수: 4,862 건 | 취소/환불 주문 수: 810 건"
])

# Slide 5: 차트 및 해석
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5)
add_header(s5, "3. 구현 결과: 차트 시각화 및 결과 해석", "DASHBOARD IMPLEMENTATION - CHARTS")
add_card(s5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "📈 Chart 1: 카테고리별 매출 현황", [
    "• Plotly Express Bar Chart (`px.bar`) 사용",
    "• 고단가 카테고리가 전체 매출 상승 견인",
    "• 주력 카테고리 중심의 마케팅 자원 집행 유효"
])
add_card(s5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "📉 Chart 2: 월별 주문 금액 추이", [
    "• Plotly Express Line Chart (`px.line`) 사용",
    "• 월별 매출 변동성 및 시즌성 패턴 확인",
    "• 성수기/비수기 맞춤 수급 계획 필요"
])

# Slide 6: 새롭게 배운 내용
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6)
add_header(s6, "4. 새롭게 배운 내용", "WHAT WE LEARNED")
add_card(s6, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "💡 핵심 학습 내용 요약", [
    "• Streamlit 반응형 UI: `@st.cache_data`를 통한 데이터 로딩 최적화 및 `st.sidebar` 필터 연동 체계 이해",
    "• Plotly Express 시각화: 인터랙티브 그래프 구현 및 `st.info` 패널을 활용한 동적 분석 결과 제공",
    "• Pandas 최신 버전에 대한 대처: Pandas 2.2.0 이후 오프셋 표기법(`resample('ME')`) 실무 적용"
])

# Slide 7: AI 코드 단위 검증
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7)
add_header(s7, "5. AI 작성 코드 단위 검증", "CODE VERIFICATION & TESTING")
add_card(s7, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "🧪 기능 단위 검증 결과 (100% PASSED)", [
    "1) Streamlit 기본 UI 검증: 타이틀, 설명문, 5종 사이드바 필터 정상 노출 (PASSED)",
    "2) KPI 연동 및 수치 정확도 검증: 필터 변경 시 8개 지표 실시간 재계산 일치 (PASSED)",
    "3) 예외 처리(Empty State) 검증: 데이터 0건 조회 시 경고 문구 출력 및 앱 다운 방지 (PASSED)",
    "4) 차트 & 결과 해석 패널 검증: Plotly 차트 및 4요소 결과 해석 노출 (PASSED)"
])

# Slide 8: 트러블슈팅
s8 = prs.slides.add_slide(blank_layout)
set_slide_background(s8)
add_header(s8, "6. 트러블슈팅 및 이슈 해결", "TROUBLESHOOTING & ISSUES")
add_card(s8, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "⚡ 주요 이슈 및 해결 과제", [
    "• 이슈 1: `df.resample('M')` 사용 시 ValueError 발생",
    "  → 해결: Pandas 2.2.0 이후 'M' 지원 중단에 따라 Month End 의미의 'ME'로 변경",
    "• 이슈 2: Plotly 라이브러리 미설치로 인한 ModuleNotFoundError",
    "  → 해결: 가상환경(`.venv`) 내 `pip install plotly` 설치 및 `requirements.txt` 명시",
    "• 이슈 3: 필터 조회 결과가 0건일 때 연산 예외 발생 위험",
    "  → 해결: `is_empty` 분기문 최상단 배치로 ZeroDivisionError 사전 차단"
])

# Slide 9: GitHub 협업
s9 = prs.slides.add_slide(blank_layout)
set_slide_background(s9)
add_header(s9, "7. GitHub 협업에서 깨달은 점", "GITHUB COLLABORATION & LESSONS")
add_card(s9, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "🌿 협업 레슨", [
    "• `requirements.txt` 의존성 관리: 팀원 간 패키지 버전 일치화로 실행 에러 방지",
    "• Feature Branch 전략 준수: 메인 브랜치 직접 커밋을 지양하고 기능 단위 브랜치 PR 작성",
    "• 코드 리뷰 문화: AI 코드를 검증 없이 병합하지 않고 로컬 테스트 후 승인 절차 진행"
])

# Slide 10: 추가 학습 계획
s10 = prs.slides.add_slide(blank_layout)
set_slide_background(s10)
add_header(s10, "8. 미해결 과제 & 추가 학습 계획", "UNRESOLVED & LEARNING PLANS")
add_card(s10, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "📚 향후 학습 방향", [
    "• Streamlit 캐싱(`@st.cache_data`) 및 세션 상태 관리(`st.session_state`) 깊이 학습",
    "• CSV 방식에서 PostgreSQL / MySQL 등 RDBMS 데이터 파이프라인 연동",
    "• 고객 RFM 세그먼테이션, 코호트 재구매율 분석 기능 시각화 고도화"
])

# Slide 11: 개선 방향
s11 = prs.slides.add_slide(blank_layout)
set_slide_background(s11)
add_header(s11, "9. 다음 프로젝트 개선 방향", "NEXT PROJECT DIRECTION")
add_card(s11, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "🚀 차기 대시보드 로드맵", [
    "1) 실시간 DB 파이프라인 구축: REST API / DB 직접 연동",
    "2) 리포트 내보내기 기능: 필터 결과 데이터 CSV 다운로드 및 PDF 보고서 생성 버튼 추가",
    "3) Multi-page App & 사용자 권한 관리: 관리자/사용자 기능 분리 대시보드 구현"
])

# Slide 12: Q&A
s12 = prs.slides.add_slide(blank_layout)
set_slide_background(s12)
tb = s12.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "감사합니다"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_MAIN
p.font.name = "Malgun Gothic"
p2 = tf.add_paragraph()
p2.text = "Q&A 및 피드백"
p2.font.size = Pt(20)
p2.font.color.rgb = PRIMARY
p2.font.name = "Malgun Gothic"

prs.save("ECommerce_Dashboard_Final_Presentation.pptx")
print("🎉 PPTX 파일 생성 완료: ECommerce_Dashboard_Final_Presentation.pptx")